"""
main.py

Task 2: Predicting Customer Buying Behaviour

Steps in this script:
1. Loads and explores the customer booking dataset (customer_booking.csv).
2. Cleans and prepares the data (handles missing values, encodes categorical variables).
3. Splits the data into training and test sets.
4. Trains a Random Forest model to predict customer bookings.
5. Evaluates the model (classification metrics, cross-validation).
6. Plots feature importance (now with improved size and horizontal bars).
7. Creates a single PowerPoint slide summarizing key findings.

Required Packages:
  pip install pandas numpy scikit-learn matplotlib seaborn python-pptx
"""

import os
import re
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns

from sklearn.model_selection import train_test_split, cross_val_score
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import classification_report, confusion_matrix, accuracy_score

from pptx import Presentation
from pptx.util import Inches

# -------------------------------------
# 1. Load and Explore the Dataset
# -------------------------------------
def load_and_explore_data(filepath):
    # Use 'latin1' encoding to handle non-UTF-8 characters
    df = pd.read_csv(filepath, encoding="latin1")
    print("Dataset Shape:", df.shape)
    print("Dataset Info:")
    print(df.info())
    print("First 5 Rows:")
    print(df.head())
    print("Summary Statistics:")
    print(df.describe())
    return df

# -------------------------------------
# 2. Clean and Prepare the Data
# -------------------------------------
def clean_and_prepare_data(df, target_column):
    # Drop rows where target is missing
    df = df.dropna(subset=[target_column])
    
    # For numeric columns, fill missing values with the median.
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    for col in numeric_cols:
        df[col] = df[col].fillna(df[col].median())
    
    # For categorical columns, fill missing values with the mode.
    categorical_cols = df.select_dtypes(include=['object']).columns.tolist()
    for col in categorical_cols:
        df[col] = df[col].fillna(df[col].mode()[0])
    
    # Convert categorical columns to dummy variables
    df_encoded = pd.get_dummies(df, columns=categorical_cols, drop_first=True)
    
    return df_encoded

# -------------------------------------
# 3. Train the Model
# -------------------------------------
def train_model(X, y):
    # Split data: 80% training, 20% testing
    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.2, random_state=42
    )
    
    model = RandomForestClassifier(n_estimators=100, random_state=42)
    model.fit(X_train, y_train)
    
    y_pred = model.predict(X_test)
    print("\nClassification Report:")
    print(classification_report(y_test, y_pred))
    print("Confusion Matrix:")
    print(confusion_matrix(y_test, y_pred))
    accuracy = accuracy_score(y_test, y_pred)
    print("Test Accuracy:", accuracy)
    
    # Cross-validation
    cv_scores = cross_val_score(model, X, y, cv=5)
    print("Cross-validation scores:", cv_scores)
    print("Mean CV Score:", cv_scores.mean())
    
    return model, X_train, X_test, y_train, y_test, accuracy, cv_scores.mean()

# -------------------------------------
# 4. Plot Feature Importance (Updated)
# -------------------------------------
def plot_feature_importance(model, feature_names, output_path):
    importances = model.feature_importances_
    indices = np.argsort(importances)[::-1]
    
    # OPTIONAL: Show only the top 15 features for readability
    top_n = 15
    indices = indices[:top_n]
    
    plt.figure(figsize=(12, 10), dpi=150)
    plt.title("Feature Importance")

    # Plot horizontal bar chart
    sns.barplot(
        x=importances[indices],
        y=[feature_names[i] for i in indices],
        orient="h"
    )
    plt.xlabel("Importance")
    plt.ylabel("Feature")
    plt.tight_layout()
    plt.savefig(output_path)
    plt.close()
    print(f"Feature importance plot saved to {output_path}")

# -------------------------------------
# 5. Create PowerPoint Slide
# -------------------------------------
def create_powerpoint_slide(total_samples, accuracy, cv_score, feature_importance_img, ppt_output):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # "Title and Content" layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add slide title
    slide.shapes.title.text = "Customer Buying Behaviour Prediction Summary"
    
    # Add key metrics textbox
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(2)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.text = (
        f"Total Samples: {total_samples}\n"
        f"Test Accuracy: {accuracy:.2f}\n"
        f"Mean CV Score: {cv_score:.2f}\n"
        "Key Predictive Variables shown in the feature importance plot."
    )
    
    # Insert the feature importance plot image
    left, top = Inches(1), Inches(4)
    slide.shapes.add_picture(feature_importance_img, left, top, width=Inches(4), height=Inches(3))
    
    prs.save(ppt_output)
    print(f"PowerPoint slide saved to {ppt_output}")

# -------------------------------------
# Main Execution
# -------------------------------------
def main():
    # Set file path for customer booking data
    data_path = "customer_booking.csv"  # Ensure this file is in your project folder
    
    # 1. Load and explore the dataset
    df = load_and_explore_data(data_path)
    
    # The target column in your dataset is "booking_complete"
    target_column = "booking_complete"
    
    # 2. Clean and prepare data
    df_prepared = clean_and_prepare_data(df, target_column)
    
    # Separate features and target variable
    X = df_prepared.drop(target_column, axis=1)
    y = df_prepared[target_column]
    
    # 3. Train the model and evaluate
    model, X_train, X_test, y_train, y_test, accuracy, mean_cv = train_model(X, y)
    
    # 4. Plot feature importance
    feature_importance_img = "feature_importance.png"
    plot_feature_importance(model, X_train.columns, feature_importance_img)
    
    # 5. Create a PowerPoint slide summarizing findings
    total_samples = df.shape[0]
    ppt_output = "Customer_Buying_Behaviour_Analysis.pptx"
    create_powerpoint_slide(total_samples, accuracy, mean_cv, feature_importance_img, ppt_output)

if __name__ == "__main__":
    main()
